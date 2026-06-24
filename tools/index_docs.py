#!/usr/bin/env python3
"""Index local documents (Word / PDF / PPT / txt / md) into docs-index.json for
MindCanvas offline 资料补充 (local file search, no internet).

Usage:
    python3 tools/index_docs.py <docs_dir> [out.json]

Defaults out.json to <docs_dir>/../docs-index.json so it sits next to the DB on the
Fly volume (point MINDCANVAS_DOCS_INDEX there, or just drop files in MINDCANVAS_DOCS_DIR).

Dependencies (install what you have docs for):
    pip install python-docx python-pptx pypdf
txt/md need nothing. Missing parsers are skipped with a warning.
"""
import sys, os, json, glob

def chunk(text, max_len=380):
    """Split text into ~paragraph-sized chunks."""
    out, buf = [], ""
    for para in [p.strip() for p in text.split("\n") if p.strip()]:
        if len(buf) + len(para) > max_len and buf:
            out.append(buf); buf = ""
        buf = (buf + " " + para).strip() if buf else para
        if len(buf) > max_len - 60:
            out.append(buf); buf = ""
    if buf:
        out.append(buf)
    return [c for c in out if len(c) >= 12]

def read_docx(path):
    from docx import Document
    paras = [p.text for p in Document(path).paragraphs if p.text.strip()]
    return [("", c) for c in chunk("\n".join(paras))]

def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path); res = []
    for i, slide in enumerate(prs.slides, 1):
        txt = "\n".join(sh.text for sh in slide.shapes if sh.has_text_frame and sh.text.strip())
        for c in chunk(txt):
            res.append((f"幻灯片{i}", c))
    return res

def read_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path); res = []
    for i, page in enumerate(reader.pages, 1):
        txt = page.extract_text() or ""
        for c in chunk(txt):
            res.append((f"第{i}页", c))
    return res

def read_text(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        return [("", c) for c in chunk(f.read())]

READERS = {".docx": read_docx, ".pptx": read_pptx, ".pdf": read_pdf,
           ".txt": read_text, ".md": read_text, ".markdown": read_text}

def main():
    if len(sys.argv) < 2:
        print(__doc__); sys.exit(1)
    docs_dir = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else os.path.join(os.path.dirname(os.path.abspath(docs_dir.rstrip("/"))), "docs-index.json")
    chunks, skipped = [], []
    for path in sorted(glob.glob(os.path.join(docs_dir, "**", "*"), recursive=True)):
        if not os.path.isfile(path):
            continue
        ext = os.path.splitext(path)[1].lower()
        reader = READERS.get(ext)
        if not reader:
            continue
        fname = os.path.basename(path)
        try:
            for loc, text in reader(path):
                chunks.append({"file": fname, "loc": loc, "text": text})
        except ImportError as e:
            skipped.append(f"{fname} (缺少库: {e.name} — pip install 后重试)")
        except Exception as e:
            skipped.append(f"{fname} ({e})")
    with open(out, "w", encoding="utf-8") as f:
        json.dump({"generatedAt": None, "chunks": chunks}, f, ensure_ascii=False, indent=0)
    print(f"✓ {len(chunks)} 段，来自 {len({c['file'] for c in chunks})} 个文件 → {out}")
    for s in skipped:
        print("  跳过:", s)

if __name__ == "__main__":
    main()
